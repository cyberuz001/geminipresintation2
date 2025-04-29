import io
import os
import re

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN

try:
    from image_scrapper import downloader
except ImportError:
    from .image_scrapper import downloader


async def generate_ppt_prompt(language, emotion_type, slide_length, topic, plan_count="5"):
    message = f"""Create an {language} language outline for a {emotion_type} slideshow presentation on the topic of {topic} which is {slide_length} slides long. 
Make sure it is {slide_length} slides long.

IMPORTANT: Include a Plan/Outline slide after the title slide with exactly {plan_count} key points that will be covered in the presentation.
Make sure the images are highly relevant to the topic and content of each slide.
Make sure the text content is well-aligned with the topic and provides valuable information.
Keep text content concise and avoid text overflow - each bullet point should be 1-2 lines maximum.

You are allowed to use the following slide types:

Slide types:
Title Slide - (Title, Subtitle)
Plan Slide - (Title, Content with {plan_count} numbered points)
Content Slide - (Title, Content)
Image Slide - (Title, Content, Image)
Thanks Slide - (Title)

Put this tag before the Title Slide: [L_TS]
Put this tag before the Plan Slide: [L_PS]
Put this tag before the Content Slide: [L_CS]
Put this tag before the Image Slide: [L_IS]
Put this tag before the Thanks Slide: [L_THS]

Put this tag after each Slide: [SLIDEBREAK]

For example:
[L_TS]
[TITLE]Mount Everest: The Highest Peak in the World[/TITLE]

[SLIDEBREAK]

[L_PS]
[TITLE]Presentation Plan[/TITLE]
[CONTENT]1. Introduction to Mount Everest
2. Geographical location and features
3. History of climbing expeditions
4. Environmental challenges
5. Future of Mount Everest tourism[/CONTENT]

[SLIDEBREAK]

[L_IS]
[TITLE]Facts about Mount Everest[/TITLE]
[CONTENT]• It is 8,848 meters (29,029 ft) high above sea level
• First successfully climbed by Sir Edmund Hillary and Tenzing Norgay on May 29, 1953
• Over 300 climbers have died attempting to scale the mountain[/CONTENT]
[IMAGE]Mount Everest snow-capped peak with climbers[/IMAGE]

[SLIDEBREAK]

Put this tag before the Title: [TITLE]
Put this tag after the Title: [/TITLE]
Put this tag before the Subtitle: [SUBTITLE]
Put this tag after the Subtitle: [/SUBTITLE]
Put this tag before the Content: [CONTENT]
Put this tag after the Content: [/CONTENT]
Put this tag before the Image: [IMAGE]
Put this tag after the Image: [/IMAGE]

Elaborate on the Content, provide as much information as possible but keep each bullet point concise.
You put a [/CONTENT] at the end of the Content.
Pay attention to the language of presentation - {language}.
Each image should be described in detail with specific keywords related to the slide content, such as "Mount Everest Summit with Climbers at Sunrise" instead of just "Mount Everest".
IMPORTANT: Always search for images in English regardless of the presentation language.
Do not write Image in Content tag.
Do not reply as if you are talking about the slideshow itself. (ex. "Include pictures here about...")
Do not include any special characters (?, !, ., :, ) in the Title.
Do not include any additional information in your response and stick to the format."""

    return message


async def generate_ppt(answer, template):
    template = os.path.join("bot", "ai_generator", "presentation_templates", f"{template}.pptx")
    root = Presentation(template)

    # """ Ref for slide types:
    # 0 -> title and subtitle
    # 1 -> title and content
    # 2 -> section header
    # 3 -> two content
    # 4 -> Comparison
    # 5 -> Title only
    # 6 -> Blank
    # 7 -> Content with caption
    # 8 -> Pic with caption
    # """

    async def delete_all_slides():
        for i in range(len(root.slides) - 1, -1, -1):
            r_id = root.slides._sldIdLst[i].rId
            root.part.drop_rel(r_id)
            del root.slides._sldIdLst[i]

    async def create_title_slide(title, subtitle):
        layout = root.slide_layouts[0]
        slide = root.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    async def create_section_header_slide(title):
        layout = root.slide_layouts[2]
        slide = root.slides.add_slide(layout)
        slide.shapes.title.text = title

    async def create_title_and_content_slide(title, content):
        layout = root.slide_layouts[1]
        slide = root.slides.add_slide(layout)
        slide.shapes.title.text = title
        
        # Format content to prevent overflow
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.text = ""  # Clear default text
        
        # Split content by lines and add as paragraphs
        lines = content.split('\n')
        for i, line in enumerate(lines):
            p = text_frame.add_paragraph()
            p.text = line.strip()
            p.font.size = Pt(18)  # Adjust font size to prevent overflow
            
            # If it's a bullet point, format it properly
            if line.strip().startswith('•'):
                p.level = 0
            elif line.strip().startswith('-'):
                p.level = 0
                
            # Limit paragraph length to prevent overflow
            if len(line) > 100:
                p.font.size = Pt(16)  # Smaller font for longer lines
            if len(line) > 150:
                p.font.size = Pt(14)  # Even smaller font for very long lines

    async def create_plan_slide(title, content):
        # Using the same layout as title and content but with special formatting
        layout = root.slide_layouts[1]
        slide = root.slides.add_slide(layout)
        slide.shapes.title.text = title
        
        # Format plan content
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.text = ""  # Clear default text
        
        # Split content by lines and add as paragraphs
        lines = content.split('\n')
        for i, line in enumerate(lines):
            p = text_frame.add_paragraph()
            p.text = line.strip()
            p.font.size = Pt(20)  # Larger font for plan points
            
            # Center align the plan points
            p.alignment = PP_ALIGN.LEFT
            
            # Add spacing between points
            if i < len(lines) - 1:
                p.space_after = Pt(12)

    async def create_title_and_content_and_image_slide(title, content, image_query):
        layout = root.slide_layouts[8]
        slide = root.slides.add_slide(layout)
        slide.shapes.title.text = title
        
        # Format content to prevent overflow
        content_shape = slide.placeholders[2]
        text_frame = content_shape.text_frame
        text_frame.text = ""  # Clear default text
        
        # Split content by lines and add as paragraphs
        lines = content.split('\n')
        for i, line in enumerate(lines):
            p = text_frame.add_paragraph()
            p.text = line.strip()
            p.font.size = Pt(16)  # Smaller font for image slides
            
            # If it's a bullet point, format it properly
            if line.strip().startswith('•'):
                p.level = 0
            elif line.strip().startswith('-'):
                p.level = 0
                
            # Limit paragraph length to prevent overflow
            if len(line) > 80:
                p.font.size = Pt(14)  # Smaller font for longer lines
            if len(line) > 120:
                p.font.size = Pt(12)  # Even smaller font for very long lines

        try:
            # Always search for images in English regardless of presentation language
            # Use more specific search terms for better image relevance
            enhanced_query = f"{title} {image_query} high quality academic"
            image_data = await downloader.download(enhanced_query, limit=1, adult_filter_off=True, timeout=15,
                                                   filter="+filterui:aspect-wide+filterui:imagesize-wallpaper+filterui:photo-photo")
            slide.placeholders[1].insert_picture(io.BytesIO(image_data))
        except Exception as e:
            print(f"Error downloading image: {e}")
            pass

    async def find_text_in_between_tags(text, start_tag, end_tag):
        start_pos = text.find(start_tag)
        end_pos = text.find(end_tag)
        result = []
        while start_pos > -1 and end_pos > -1:
            text_between_tags = text[start_pos + len(start_tag):end_pos]
            result.append(text_between_tags)
            start_pos = text.find(start_tag, end_pos + len(end_tag))
            end_pos = text.find(end_tag, start_pos)
        res1 = "".join(result)
        res2 = re.sub(r"\[IMAGE\].*?\[/IMAGE\]", '', res1)
        if len(result) > 0:
            return res2
        else:
            return ""

    async def search_for_slide_type(text):
        tags = ["[L_TS]", "[L_PS]", "[L_CS]", "[L_IS]", "[L_THS]"]
        found_text = next((s for s in tags if s in text), None)
        return found_text

    async def parse_response(reply):
        list_of_slides = reply.split("[SLIDEBREAK]")
        for slide in list_of_slides:
            slide_type = await search_for_slide_type(slide)
            match slide_type:
                case ("[L_TS]"):
                    await create_title_slide(await find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]"),
                                             await find_text_in_between_tags(str(slide), "[SUBTITLE]", "[/SUBTITLE]"))
                case ("[L_PS]"):
                    await create_plan_slide("".join(await find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")),
                                            "".join(await find_text_in_between_tags(str(slide), "[CONTENT]", "[/CONTENT]")))
                case ("[L_CS]"):
                    await create_title_and_content_slide("".join(await find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")),
                                                         "".join(await find_text_in_between_tags(str(slide), "[CONTENT]", "[/CONTENT]")))
                case ("[L_IS]"):
                    await create_title_and_content_and_image_slide("".join(await find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")),
                                                                   "".join(await find_text_in_between_tags(str(slide), "[CONTENT]", "[/CONTENT]")),
                                                                   "".join(await find_text_in_between_tags(str(slide), "[IMAGE]", "[/IMAGE]")))
                case ("[L_THS]"):
                    await create_section_header_slide("".join(await find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")))

    async def find_title():
        return root.slides[0].shapes.title.text

    await delete_all_slides()
    await parse_response(answer)
    buffer = io.BytesIO()
    root.save(buffer)
    pptx_bytes = buffer.getvalue()
    pptx_title = f"{await find_title()}.pptx"
    print(f"done {pptx_title}")

    return pptx_bytes, pptx_title
